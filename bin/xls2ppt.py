from pptx import Presentation, util #requires python-pptx
from pptx.enum.dml import MSO_THEME_COLOR
from openpyxl import Workbook, load_workbook #requires openpyxl
import re
import copy, six, os, sys

# From https://github.com/scanny/python-pptx/issues/132
def _get_blank_slide_layout(pres):
    layout_items_count = [len(layout.placeholders) for layout in pres.slide_layouts]
    min_items = min(layout_items_count)
    blank_layout_id = layout_items_count.index(min_items)
    #return pres.slide_layouts[blank_layout_id]
    return pres.slide_layouts[4] #manually set empty slide


def list_layout_names(prs):
    """
    Lists the available index and layout-names to determine the empty, to use, slide layout
    :param prs: the presentation object
    :return:
    """
    for ii, l in enumerate(prs.slide_layouts):
        print("prs.slide_layouts[{:d}]: {:s}".format(ii, l.name))

def duplicate_slide(pres, index):
    """Duplicate the slide with the given index in pres.

    Adds slide to the end of the presentation"""
    source = pres.slides[index]

    blank_slide_layout = _get_blank_slide_layout(pres)
    dest = pres.slides.add_slide(blank_slide_layout)

    for shp in source.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

    # Remove empty placeholders
    for shp in dest.shapes:
        if shp.is_placeholder:
            el = shp.element
            el.getparent().remove(el)

    try:
        for key, value in enumerate(source.rels):
            # Make sure we don't copy a notesSlide relation as that won't exist
            if not "notesSlide" in value.reltype:
                dest.rels.add_relationship(value.reltype, value._target, value.rId)
    except Exception:
        # dont do anything
        print("No relationships detected")

    return dest


def query_yes_no(question, default="yes"):
    """Ask a yes/no question via raw_input() and return their answer.

    "question" is a string that is presented to the user.
    "default" is the presumed answer if the user just hits <Enter>.
        It must be "yes" (the default), "no" or None (meaning
        an answer is required of the user).

    The "answer" return value is True for "yes" or False for "no".
    """
    valid = {"yes": True, "y": True, "ye": True,
             "no": False, "n": False}
    if default is None:
        prompt = " [y/n] "
    elif default == "yes":
        prompt = " [Y/n] "
    elif default == "no":
        prompt = " [y/N] "
    else:
        raise ValueError("invalid default answer: '%s'" % default)

    while True:
        sys.stdout.write(question + prompt)
        choice  = input().lower()
        if default is not None and choice == '':
            return valid[default]
        elif choice in valid:
            return valid[choice]
        else:
            sys.stdout.write("Please respond with 'yes' or 'no' "
                             "(or 'y' or 'n').\n")

def add_run(p, text, ph):
    run = p.add_run()
    # Add text to current run
    run.text = str(text)
    # Format the current run
    font = run.font
    font.name = ph['font'] if ph['font'] else font.name
    font.size = util.Pt(ph['fontsize']) if ph['fontsize'] else font.size
    font.bold = ph['bold']
    font.italic = None
    if ph['color_idx']:
        font.color.theme_color = ph['color_idx']
    return run

# ----------------------------------------------------------------
# parenDir =       r'C:\temp'
# pptTemplateFile = os.path.join(parenDir, 'Template.pptx')
# xlsDataFile =     os.path.join(parenDir, 'Content.xlsx')
# pptOutFile =      os.path.join(parenDir, 'Output.pptx')





#regexp_pat = re.compile("@([A-Z]+)(([0-9]+)(:([0-9]+))?)?")
regexp_pat = re.compile("@(.+)")

# Define pattern to identify bold font
pat_bold = re.compile('\*[^\*]+\*')


prs = Presentation(pptTemplateFile)
wb = load_workbook(xlsDataFile, data_only=True)

slide = prs.slides[0]

# grab the active worksheet
ws = wb.active
max_row = ws.max_row
max_col = ws.max_column
# identify the header rows
header_row_until = []
for rowIdx in range(1, max_row+1):
    if ws['A{:d}'.format(rowIdx)].value is None:
        header_row_until = rowIdx-1
        print("Found last header row: {:d}".format(rowIdx))
        break
if header_row_until is []:
    raise Exception('No end of header found')

# generate list of placeholder tags
placeholder = dict()
for rowIdx in range(1, header_row_until+1):
    if ws.cell(row=rowIdx, column=1).value == "Placeholder":
        for colIdx in range(2,max_col+1):
            ph = regexp_pat.search(ws.cell(column=colIdx, row=rowIdx).value)
            if ph:
                placeholder[ph.group(1)] = dict()
                placeholder[ph.group(1)]['colIdx'] = colIdx

# generate list of accent colors
MSO_THEME_COLOR_list = [MSO_THEME_COLOR.ACCENT_1,
                                    MSO_THEME_COLOR.ACCENT_2,
                                    MSO_THEME_COLOR.ACCENT_3,
                                    MSO_THEME_COLOR.ACCENT_4,
                                    MSO_THEME_COLOR.ACCENT_5,
                                    MSO_THEME_COLOR.ACCENT_6]
# generate font-list
for rowIdx in range(1, header_row_until+1):
    if ws.cell(row=rowIdx, column=1).value == "Font":
        for ph in placeholder:
            placeholder[ph]['font'] = ws.cell(row=rowIdx, column=placeholder[ph]['colIdx']).value

    if ws.cell(row=rowIdx, column=1).value == "Fontsize":
        for ph in placeholder:
            placeholder[ph]['fontsize'] = ws.cell(row=rowIdx, column=placeholder[ph]['colIdx']).value
            if placeholder[ph]['fontsize'] is None:
                placeholder[ph]['fontsize'] = None
            else:
                placeholder[ph]['fontsize'] = int(placeholder[ph]['fontsize'])

    if ws.cell(row=rowIdx, column=1).value == "Bold":
        for ph in placeholder:
            placeholder[ph]['bold'] = True if ws.cell(row=rowIdx, column=placeholder[ph]['colIdx']).value is not None else False

    if ws.cell(row=rowIdx, column=1).value == "ThemeColorIndex":
        for ph in placeholder:
            placeholder[ph]['color_idx'] = ws.cell(row=rowIdx, column=placeholder[ph]['colIdx']).value
            if placeholder[ph]['color_idx'] is not None and placeholder[ph]['color_idx'] != "":
                placeholder[ph]['color_idx'] = MSO_THEME_COLOR_list[int(placeholder[ph]['color_idx'])-1]
            else:
                placeholder[ph]['color_idx'] = None




duplicateSlides = True
print("Using rows with content as row Counter. max_row: {}".format(max_row))

# Loop through Rows until last one hit
for rowIdx in range(header_row_until+1,max_row+1):
    # Duplicate First Slide with all the Placeholders if Slides should be duplicated
    if duplicateSlides:
        slide = duplicate_slide(prs, 0)
        # select last slide
        slide = prs.slides[prs.slides.__len__()-1]
        # Setting a name for a slide does not work :( - you would see the name when setting a link within the same ppt
        # slide.name = 'foo {:d}'.format(prs.slides.__len__()+1)
    else:
        # select slide for insertation
        slide = prs.slides[rowIdx]
    # Loop through the shapes on the slide and look for the Identifiers
    for shape in slide.shapes:
        # Check if shape has "@A" identifier
        if shape.has_text_frame:
            m = regexp_pat.search(shape.text)
            if m:
                # cell = m.group(1)
                # rowStart = m.group(3)
                # rowEnd = m.group(5)
                tf = shape.text_frame
                mm = regexp_pat.search(tf.text)
                if mm and mm.group(1) in placeholder.keys():
                    ph = placeholder[mm.group(1)]
                    replace_text = ws.cell(row=rowIdx, column=ph['colIdx']).value
                    print("Inserting in Shape {} in paragraph value {}".format(shape.text, replace_text))
                    # Clear current content of text_frame
                    tf.clear()
                    if replace_text is not None:
                        if "*" in replace_text:
                            print('debug')

                        # Divide Text in bold and not-bold parts
                        not_bold_replace_text = pat_bold.split(replace_text)
                        bold_replace_text = pat_bold.findall(replace_text)

                        p = tf.paragraphs[0]

                        for ii_bold in range(0, len(bold_replace_text)):
                            # Add not_bold_text
                            add_run(p, not_bold_replace_text[ii_bold], ph)

                            # Add bold_text
                            ph_bold = ph.copy()
                            ph_bold['bold'] = True
                            add_run(p, bold_replace_text[ii_bold][1:-1], ph_bold)

                        # Add remaining not_bold_text
                        add_run(p, not_bold_replace_text[-1], ph)



                    # name the shape - for later import of modified slides
                    shape.name = '@{:s}({:d},{:d})'.format(mm.group(1), rowIdx, ph['colIdx'])

                # if re.search("\n",str(ws[m.group(1)+str(rowIdx)].value)):
                #     # Seems like a multi-line content. Better put it into a bullet-list
                #     p = tf
                #     for paragraph in re.split("\n", str(ws[m.group(1)+str(rowIdx)].value)):
                #         p.text = paragraph
                #         p = tf.add_paragraph()
        else:
            print("Possibly a Group. They are note yet supported")

if True: #query_yes_no('Datei "{}" schon vorhanden. Überschreiben?'.format(os.path.basename(pptOutFile)),default="yes"):
    prs.save(pptOutFile)
    os.startfile(pptOutFile)
else:
    print("Keine Datei gespeichert!")